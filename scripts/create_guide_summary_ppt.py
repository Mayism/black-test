from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "GUIDE_SUMMARY.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

FONT = "Microsoft YaHei"
TITLE_FONT = "Microsoft YaHei UI"

NAVY = RGBColor(28, 43, 64)
BLUE = RGBColor(51, 111, 197)
TEAL = RGBColor(34, 150, 137)
GREEN = RGBColor(69, 160, 92)
AMBER = RGBColor(219, 146, 45)
RED = RGBColor(204, 74, 74)
GRAY = RGBColor(91, 104, 124)
LIGHT_BG = RGBColor(247, 249, 252)
BORDER = RGBColor(217, 224, 235)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(30, 36, 45)


def set_bg(slide, color=LIGHT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def text_box(slide, x, y, w, h, text, size=16, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.06)
    frame.margin_right = Inches(0.06)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def title(slide, main, sub=None):
    text_box(slide, 0.55, 0.34, 9.4, 0.46, main, 24, NAVY, True)
    if sub:
        text_box(slide, 0.58, 0.83, 10.8, 0.28, sub, 10.5, GRAY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(1.15), Inches(1.2), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def footer(slide, page):
    text_box(slide, 0.58, 7.08, 5.2, 0.22, "PRD 到 HarmonyRun 测试套件工程说明", 8.5, GRAY)
    text_box(slide, 11.75, 7.08, 1.0, 0.22, f"{page:02d}", 9, GRAY, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, head, lines, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    text_box(slide, x + 0.2, y + 0.14, w - 0.36, 0.32, head, 15, NAVY, True)
    box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.58), Inches(w - 0.36), Inches(h - 0.68))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(11.3)
        p.font.color.rgb = BLACK
        p.space_after = Pt(5)


def code_box(slide, x, y, w, h, code, size=8.7):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.adjustments[0] = 0.04
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(35, 42, 55)
    shape.line.color.rgb = RGBColor(68, 78, 96)
    box = slide.shapes.add_textbox(Inches(x + 0.14), Inches(y + 0.12), Inches(w - 0.28), Inches(h - 0.24))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for i, line in enumerate(code.splitlines()):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(235, 240, 247)
        p.space_after = Pt(0)
    return shape


def stat(slide, x, y, w, h, number, label, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER
    text_box(slide, x, y + 0.18, w, 0.52, str(number), 27, color, True, PP_ALIGN.CENTER)
    text_box(slide, x + 0.08, y + 0.82, w - 0.16, 0.28, label, 10.5, GRAY, False, PP_ALIGN.CENTER)


def table(slide, x, y, w, h, data, col_widths, font_size=10):
    shape = slide.shapes.add_table(len(data), len(data[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = shape.table
    for idx, cw in enumerate(col_widths):
        tbl.columns[idx].width = Inches(cw)
    for r, row in enumerate(data):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = value
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r == 0 else (WHITE if r % 2 else RGBColor(243, 247, 252))
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT
                p.font.size = Pt(font_size + 0.5 if r == 0 else font_size)
                p.font.bold = r == 0
                p.font.color.rgb = WHITE if r == 0 else BLACK
                p.alignment = PP_ALIGN.CENTER if c >= 2 else PP_ALIGN.LEFT
    return shape


def arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = GRAY
    line.line.width = Pt(1.4)
    line.line.end_arrowhead = True


def new_slide(page, main, sub=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title(s, main, sub)
    footer(s, page)
    return s


# 1
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, RGBColor(244, 248, 251))
text_box(s, 0.75, 0.72, 8.7, 0.42, "工程汇报", 16, TEAL, True)
text_box(s, 0.72, 1.25, 9.6, 1.1, "PRD 到 HarmonyRun\n测试套件生成与执行", 36, NAVY, True)
text_box(s, 0.78, 2.65, 8.8, 0.52, "功能、架构、测试套 JSON、生成原则、测试结果与演进方向", 18, GRAY)
card(s, 0.78, 4.45, 3.6, 1.18, "当前工程目录", r"E:\black\prds", TEAL)
card(s, 4.68, 4.45, 3.6, 1.18, "汇报依据", "GUIDE.md + result 执行报告", BLUE)
card(s, 8.58, 4.45, 3.6, 1.18, "页数", "9 页精简版", AMBER)
footer(s, 1)

# 2
s = new_slide(2, "工程功能", "把 PRD Markdown 转换为 HarmonyRun 可执行、可落盘、可复盘的黑盒 UI 测试资产")
stat(s, 0.72, 1.35, 2.25, 1.18, "4", "每个 PRD 生成四件套", BLUE)
stat(s, 3.25, 1.35, 2.25, 1.18, "3", "支持 PRD 分类", TEAL)
stat(s, 5.78, 1.35, 2.25, 1.18, "1", "自动调用执行链路", GREEN)
stat(s, 8.31, 1.35, 2.25, 1.18, "N", "批量处理与归档", AMBER)
card(s, 0.72, 3.05, 3.75, 2.55, "输入", ["读取 prd/ 下 Markdown", "支持单 PRD、分类、all", "CLI 参数覆盖 config.yaml"], BLUE)
card(s, 4.78, 3.05, 3.75, 2.55, "生成", ["调用 OpenCode HTTP API", "使用 harmonyrun-testcase-gen skill", "生成 JSON、MD、app-card、agent-prompt"], TEAL)
card(s, 8.84, 3.05, 3.75, 2.55, "执行", ["调用 harmonyrun test", "真机黑盒 UI 驱动", "报告、截图、轨迹写入 result/"], GREEN)

# 3
s = new_slide(3, "架构与运行链路", "脚本负责编排，skill 负责生成，HarmonyRun 负责执行")
flow = [("PRD 输入", "文件/分类/all"), ("脚本编排", "bat/Python"), ("OpenCode", "session+prompt"), ("Skill 生成", "测试四件套"), ("JSON 校验", "suite+test_cases"), ("HarmonyRun", "真机执行"), ("结果落盘", "report/trace")]
for i, (head, body) in enumerate(flow):
    x = 0.58 + i * 1.82
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.9), Inches(1.48), Inches(1.18))
    box.adjustments[0] = 0.08
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = [BLUE, TEAL, AMBER, GREEN, BLUE, RED, NAVY][i]
    box.line.width = Pt(1.5)
    text_box(s, x + 0.04, 2.13, 1.4, 0.26, head, 12.2, NAVY, True, PP_ALIGN.CENTER)
    text_box(s, x + 0.04, 2.53, 1.4, 0.22, body, 8.8, GRAY, False, PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        arrow(s, x + 1.48, 2.49, x + 1.76, 2.49)
card(s, 0.75, 4.05, 5.65, 1.55, "主脚本职责", "读取配置、解析 PRD 范围、启动/复用 OpenCode、监听生成过程、校验 JSON、按配置调用 HarmonyRun。", BLUE)
card(s, 6.95, 4.05, 5.65, 1.55, "HarmonyRun 边界", "执行真实 UI 上可观察的页面跳转、文本、按钮、弹窗和布局；不负责后台造数、数据库校验或真实支付闭环。", TEAL)

# 4
s = new_slide(4, "测试套 JSON 格式", "JSON 是 HarmonyRun 执行入口；辅助文件负责描述应用结构和 agent 行为约束")
card(s, 0.72, 1.35, 2.75, 4.55, "根节点", ["suite：套件元信息", "test_cases：用例数组", "不额外扩展非 schema 字段"], BLUE)
card(s, 3.65, 1.35, 2.75, 4.55, "suite", ["id / name：套件标识", "app_package：目标应用包名", "app_card：应用结构卡片", "agent_prompt：执行约束"], TEAL)
card(s, 6.58, 1.35, 2.75, 4.55, "test_cases[]", ["id / title：用例标识与标题", "level：MUST / SHOULD 等级", "preconditions：前置条件", "test_steps：操作步骤", "expected_result：最终预期"], GREEN)
card(s, 9.51, 1.35, 2.75, 4.55, "test_steps[]", ["step：步骤序号", "action：单一 UI 动作", "checkpoint：可观察验证点", "每步只做一个动作"], AMBER)

# 5
s = new_slide(5, "JSON 示例", "短链路、UI 可观察、通过 file: 引用 app-card 和 agent-prompt")
example = """{
  "suite": {
    "id": "tourist_test_suite",
    "name": "常用旅客管理测试套件",
    "app_package": "com.huawei.rentandbuy",
    "app_card": "file:./tourist-app-card.md",
    "agent_prompt": "file:./tourist-agent-prompt.md"
  },
  "test_cases": [{
    "id": "MUST-02",
    "title": "进入添加常用旅客编辑页",
    "level": "MUST",
    "preconditions": ["应用已启动", "网络连接正常"],
    "test_steps": [
      {
        "step": 1,
        "action": "点击【我的】页的【常用旅客】入口",
        "checkpoint": "进入列表页，展示【添加】入口"
      },
      {
        "step": 2,
        "action": "点击【添加】",
        "checkpoint": "进入编辑页，核心表单项可见"
      }
    ],
    "expected_result": "可进入添加常用旅客编辑页"
  }]
}"""
code_box(s, 0.65, 1.32, 7.15, 5.45, example, 7.0)
card(s, 8.1, 1.35, 4.25, 1.38, "关键点", ["HarmonyRun 能直接解析执行", "checkpoint 都是 UI 可观察结果", "失败时能定位到具体步骤"], BLUE)
card(s, 8.1, 3.12, 4.25, 1.38, "不写进 JSON", ["不新增 device_type、breakpoint、assertions、depends_on 等字段", "断点和适配规则放进文本字段或辅助文件"], RED)
card(s, 8.1, 4.9, 4.25, 1.38, "辅助文件", ["app-card 描述页面结构、入口和 fixture 边界", "agent-prompt 约束执行方式、收尾和失败判定"], TEAL)

# 6
s = new_slide(6, "用例生成原则", "先判断 PRD 类型，再选择生成策略；默认生成可执行、短链路、UI 可观察的用例")
table(s, 0.62, 1.35, 12.1, 2.85, [
    ["PRD 类型", "判定信号", "生成策略"],
    ["元服务/业务功能开发", "页面、入口、表单、业务流程、账号、权限、搜索、状态流转", "短链路、独立可运行、UI 可观察；L0 主链路优先"],
    ["一多适配", "多设备、响应式、断点、横竖屏、自由窗口、分栏、导航形态", "按 widthBp/heightBp 组织，检查布局、可读性、交互可达性"],
    ["极简一多适配", "一句话：实现 XX 元服务的一多适配", "只生成可达页面适配基线巡检，结论保守表达"],
    ["混合型", "业务功能 + 明确适配要求", "先覆盖 L0 主链路，再补 PRD 明确的一多适配点"],
], [2.1, 4.55, 5.45], 9.4)
card(s, 0.72, 4.55, 3.75, 1.3, "业务功能用例", ["3-5 步短链路", "checkpoint 必须 UI 可观察", "fixture 缺口不硬写 JSON"], BLUE)
card(s, 4.78, 4.55, 3.75, 1.3, "一多适配用例", ["优先 sm/md/lg 断点", "通用规则看明显错乱", "定制规则看栅格、分栏、导航切换"], TEAL)
card(s, 8.84, 4.55, 3.75, 1.3, "结论边界", ["可证明到达页面无明显问题", "不能证明所有页面都适配或功能完整正常"], AMBER)

# 7
s = new_slide(7, "测试结果：full-generation", "列出当前 result/full-generation 下每个测试套的用例结果")
table(s, 0.45, 1.23, 12.45, 5.62, [
    ["测试套", "用例", "测试项", "结果", "备注"],
    ["household", "MUST-01", "首页基础结构展示", "通过", "首页核心结构展示正常"],
    ["household", "MUST-02", "搜索服务并查看搜索结果", "失败", "final-chance 仍输出 type_at"],
    ["household", "MUST-03", "从首页查看服务详情页", "失败", "final-chance 仍输出 click"],
    ["household", "SHOULD-01", "切换城市后首页城市名称更新", "失败", "未进入城市选择页"],
    ["household", "SHOULD-02", "全部服务页面分类导航切换", "通过", "分类联动正常"],
    ["household", "SHOULD-03", "未关联账号订单引导", "失败", "账号状态不满足"],
    ["scenic", "MUST-01", "首页TabBar与轮播展示", "通过", "展示正常"],
    ["scenic", "MUST-02", "门票Tab浏览门票列表", "通过", "列表正常"],
    ["scenic", "MUST-03", "导览Tab地图景点展示", "通过", "地图与景点入口正常"],
    ["scenic", "MUST-04", "我的Tab用户信息展示", "通过", "入口展示正常"],
    ["scenic", "SHOULD-01", "首页热门景点模块展示", "通过", "卡片展示正常"],
    ["scenic", "SHOULD-02", "首页攻略游记模块展示", "通过", "条目展示正常"],
    ["scenic", "SHOULD-03", "门票列表点击查看详情", "通过", "可进入详情页"],
    ["scenic", "SHOULD-04", "未关联账号订单引导", "失败", "未出现账号关联引导"],
], [1.35, 1.15, 3.65, 0.85, 5.45], 6.7)

# 8
s = new_slide(8, "测试结果：requirement", "列出当前 result/requirement 下每个测试套的用例结果")
table(s, 0.45, 1.28, 12.45, 5.45, [
    ["测试套", "用例", "测试项", "结果", "备注"],
    ["hotscenic", "MUST-01", "首页热门景点模块展示", "通过", "模块、更多入口和卡片正常"],
    ["hotscenic", "MUST-02", "点击热门景点卡片进入详情页", "通过", "详情页核心信息正常"],
    ["hotscenic", "SHOULD-01", "点击更多进入景点列表页", "通过", "列表页结构正常"],
    ["hotscenic", "SHOULD-02", "景点列表页点击景点进入详情页", "通过", "可进入详情页"],
    ["hotscenic", "SHOULD-03", "详情页语音讲解功能", "失败", "未观察到播放态或权限弹窗"],
    ["hotscenic", "SHOULD-04", "详情页地图导航功能", "通过", "可触发导航交互"],
    ["tourist", "MUST-01", "常用旅客列表页浏览", "通过", "列表结构和添加入口正常"],
    ["tourist", "MUST-02", "进入添加常用游客编辑页", "通过", "表单项展示正常"],
    ["tourist", "MUST-03", "添加常用游客并保存", "通过", "保存流程通过"],
    ["tourist", "SHOULD-01", "必填字段为空校验", "通过", "提示完整信息"],
    ["tourist", "SHOULD-02", "游客手机号格式校验", "通过", "提示手机号格式错误"],
    ["tourist", "SHOULD-03", "游客身份证号格式校验", "通过", "提示身份证号格式错误"],
], [1.35, 1.15, 3.85, 0.85, 5.25], 7.3)

# 9
s = new_slide(9, "演进和优化方向", "围绕生成质量、执行稳定性和结论可信度持续收敛")
card(s, 0.72, 1.35, 3.75, 2.05, "用例生成质量", ["明确组件模板用例测试范围", "区分 mock 与非 mock 场景", "避免不可控真实支付/密码等流程"], BLUE)
card(s, 4.78, 1.35, 3.75, 2.05, "结果反哺 skill", ["批量执行评测平台用例", "分析失败与真实效果不符的案例", "针对输入法遮挡、收尾失败等反向优化"], TEAL)
card(s, 8.84, 1.35, 3.75, 2.05, "一多适配标准", ["完善响应式布局之外的评价标准", "补充断点、组件、自由窗口规则", "沉淀通用规则与最佳实践规则"], AMBER)
card(s, 0.72, 4.05, 5.78, 1.62, "HarmonyRun 稳定性", "跟踪稳定版本发布进展，对比不同版本执行效果，持续归因 agent 收尾、识别和环境问题。", GREEN)
card(s, 6.82, 4.05, 5.78, 1.62, "云手机接入", "跟进云手机测试对接，提升设备环境可复制性，降低本地真机账号态和设备态差异。", RED)

prs.core_properties.title = "PRD 到 HarmonyRun 测试套件工程汇报"
prs.core_properties.subject = "GUIDE.md 汇总 PPT"
prs.core_properties.author = "Codex"
prs.save(OUT)
print(OUT)
